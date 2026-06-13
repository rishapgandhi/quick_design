"""PPTX generation via python-pptx. Gemini produces structured JSON, we render to .pptx."""

import json
import os
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = Path(os.environ.get("OUTPUT_DIR", Path(__file__).parent.parent.parent / "output"))


PPTX_PROMPT = """Generate a presentation as a JSON array of slides. Each slide is an object with:
- "title": string (slide heading)
- "content": array of bullet strings (3-5 bullets per slide)
- "notes": string (optional speaker notes)

The presentation should be professional, concise, and impactful.
Use short punchy bullets (max 10 words each).
First slide is the title slide (title = product name, content = [tagline]).
Last slide is the CTA/closing slide.

Return ONLY valid JSON. No markdown, no explanation. Example format:
[
  {"title": "ProductName", "content": ["Your compelling tagline here"], "notes": "Introduction"},
  {"title": "The Problem", "content": ["Pain point 1", "Pain point 2", "Pain point 3"], "notes": ""},
  {"title": "Get Started", "content": ["website.com", "email@company.com"], "notes": "Closing"}
]"""


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def create_pptx(slides_data: list[dict], title: str, colors: dict = None) -> Path:
    """Create a .pptx file from structured slide data."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    if colors is None:
        colors = {
            "bg": "#1a1a2e",
            "title": "#ffffff",
            "body": "#cccccc",
            "accent": "#e94560",
        }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(colors["bg"])

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(44) if i == 0 else Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(colors["title"])
        p.alignment = PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER

        # Content bullets
        content = slide_data.get("content", [])
        if content:
            body_top = Inches(3.5) if i == 0 else Inches(3.8)
            body_box = slide.shapes.add_textbox(Inches(1), body_top, Inches(11), Inches(3.5))
            tf = body_box.text_frame
            tf.word_wrap = True

            for j, bullet in enumerate(content):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(24) if i == 0 else Pt(20)
                p.font.color.rgb = hex_to_rgb(colors["body"])
                p.space_after = Pt(12)
                p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12), Inches(6.8), Inches(1), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i + 1}/{len(slides_data)}"
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgb("#666666")
        p.alignment = PP_ALIGN.RIGHT

        # Speaker notes
        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # Save
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    safe_title = "".join(c if c.isalnum() or c in "-_" else "-" for c in title[:30]).strip("-").lower()
    filename = f"{safe_title}-deck-{ts}.pptx"
    output_path = OUTPUT_DIR / filename
    prs.save(str(output_path))
    return output_path


def get_color_scheme(design_system_id: str) -> dict:
    """Extract colors from a design system for PPTX use."""
    # Dark themes
    dark_systems = {
        "dark-modern": {"bg": "#0d0d0d", "title": "#ffffff", "body": "#a0a0a0", "accent": "#ff6b35"},
        "linear": {"bg": "#08090a", "title": "#f7f8f8", "body": "#d0d6e0", "accent": "#5e6ad2"},
        "claude": {"bg": "#1a1110", "title": "#ffffff", "body": "#ccbcb5", "accent": "#d97757"},
    }
    # Light themes
    light_systems = {
        "default": {"bg": "#1e1e2e", "title": "#ffffff", "body": "#cccccc", "accent": "#2F6FEB"},
        "stripe": {"bg": "#0a2540", "title": "#ffffff", "body": "#d0d8e0", "accent": "#635bff"},
        "vercel": {"bg": "#000000", "title": "#ffffff", "body": "#999999", "accent": "#ffffff"},
        "notion": {"bg": "#191919", "title": "#ffffff", "body": "#b0b0b0", "accent": "#2eaadc"},
    }
    all_schemes = {**dark_systems, **light_systems}
    return all_schemes.get(design_system_id, all_schemes["default"])
